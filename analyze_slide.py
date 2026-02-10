from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import json, os

pptx_path = r"d:\Dev\AI테스트용\pptx2web\중등 2-1 2단원_0209(최종)_익히기합본.pptx"
prs = Presentation(pptx_path)

# Slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height
print(f"Slide size: {slide_width} x {slide_height} EMU")
print(f"Slide size: {slide_width/914400:.2f} x {slide_height/914400:.2f} inches")
print(f"Slide size: {slide_width/12700:.0f} x {slide_height/12700:.0f} pt")

# Get slide 2 (index 1)
slide = prs.slides[1]
print(f"\n=== SLIDE 2 ===")
print(f"Slide layout: {slide.slide_layout.name}")

# Check slide background
bg = slide.background
if bg.fill.type is not None:
    print(f"Background fill type: {bg.fill.type}")

print(f"\nTotal shapes: {len(slide.shapes)}")
print("=" * 80)

for i, shape in enumerate(slide.shapes):
    print(f"\n--- Shape {i} ---")
    print(f"  Name: {shape.shape_id} / {shape.name}")
    print(f"  Shape type: {shape.shape_type}")
    print(f"  Position: left={shape.left}, top={shape.top}")
    print(f"  Size: width={shape.width}, height={shape.height}")
    print(f"  Position (pt): left={shape.left/12700:.1f}, top={shape.top/12700:.1f}")
    print(f"  Size (pt): width={shape.width/12700:.1f}, height={shape.height/12700:.1f}")
    print(f"  Rotation: {shape.rotation}")
    
    if hasattr(shape, "text") and shape.text:
        print(f"  Text: '{shape.text[:100]}'")
    
    if shape.has_text_frame:
        tf = shape.text_frame
        for pi, para in enumerate(tf.paragraphs):
            print(f"  Paragraph {pi}: alignment={para.alignment}")
            for ri, run in enumerate(para.runs):
                font = run.font
                print(f"    Run {ri}: '{run.text[:80]}'")
                print(f"      Font: name={font.name}, size={font.size}, bold={font.bold}, italic={font.italic}")
                if font.color and font.color.type is not None:
                    try:
                        print(f"      Color: {font.color.rgb}")
                    except:
                        print(f"      Color: theme_color={font.color.theme_color}")
    
    # Check for fill
    if hasattr(shape, 'fill'):
        try:
            fill = shape.fill
            if fill.type is not None:
                print(f"  Fill type: {fill.type}")
                if fill.type == 1:  # Solid
                    try:
                        print(f"  Fill color: {fill.fore_color.rgb}")
                    except:
                        print(f"  Fill color: theme")
        except:
            pass
    
    # Check for image
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        print(f"  Image: content_type={image.content_type}, size={len(image.blob)} bytes")
    
    # Check for group
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"  Group with {len(shape.shapes)} sub-shapes")
        for j, sub in enumerate(shape.shapes):
            print(f"    Sub-shape {j}: {sub.name}, type={sub.shape_type}")
            print(f"      Pos: left={sub.left}, top={sub.top}, w={sub.width}, h={sub.height}")
            if hasattr(sub, "text") and sub.text:
                print(f"      Text: '{sub.text[:60]}'")
    
    # Check for table
    if shape.has_table:
        table = shape.table
        print(f"  Table: {table.rows.__len__()} rows x {len(table.columns)} cols")
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text.strip():
                    print(f"    Cell[{row_idx},{col_idx}]: '{cell.text[:50]}'")

    # Line/connector
    if hasattr(shape, 'line'):
        try:
            line = shape.line
            if line.fill.type is not None:
                print(f"  Line: width={line.width}")
                try:
                    print(f"  Line color: {line.color.rgb}")
                except:
                    pass
        except:
            pass
