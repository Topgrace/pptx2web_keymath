from pptx import Presentation
from pptx.oxml.ns import qn
import base64

pptx_path = r"d:\Dev\AI테스트용\pptx2web\중등 2-1 2단원_0209(최종)_익히기합본.pptx"
prs = Presentation(pptx_path)
slide = prs.slides[1]

# Extract all shapes' fill and text details for the "2^3" area
# Shape 4: white rounded rect background at (39,169) 626x129 - this is the "2³" container
# Shape 5: "2" large number at (186, 186) - base number
# Shape 13: "3" superscript at (248, 174) - exponent
# Shape 14: "지수" label at (425, 186)
# Shape 15: "밑" label at (425, 246)
# Shape 16: arrow line from 2³ to "지수" label at (296, 204) width 129
# Shape 17: arrow line from 2 to "밑" label at (260, 264) width 165

# Check the "2×2×2" part under the 2³
# Actually let me re-examine the middle area shapes
print("=== Middle section shapes (100 < y < 350) ===")
for i, shape in enumerate(slide.shapes):
    y_pt = shape.top / 12700
    if 100 < y_pt < 350:
        txt = ""
        if hasattr(shape, 'text') and shape.text:
            txt = shape.text
        print(f"  Shape {i}: ({shape.left/12700:.0f},{y_pt:.0f}) {shape.width/12700:.0f}x{shape.height/12700:.0f} name='{shape.name}' text='{txt}'")

# Let me also check the group with "a" (it's actually the math equation a^1 = a)
# Shape 20: group at (407,380) containing math "a" and superscript "1"
# Shape 21: "=" at (452, 392)
# So: a¹ = a ... meaning exponent 1 can be omitted

# Formula at bottom: aⁿ = a × a × ... × a with bracket "n번"
# Let me check what shapes form "a × a × ... × a"
# and aⁿ

# Get the bottom formula area (y > 370)
print("\n=== Full bottom formula area ===")
for i, shape in enumerate(slide.shapes):
    y_pt = shape.top / 12700
    x_pt = shape.left / 12700
    w_pt = shape.width / 12700
    h_pt = shape.height / 12700
    if y_pt > 370:
        txt = ""
        if hasattr(shape, 'text') and shape.text:
            txt = shape.text
        if shape.shape_type == 6:  # GROUP
            for sub in shape.shapes:
                stxt = sub.text if hasattr(sub, 'text') and sub.text else ""
                print(f"    Sub: ({sub.left/12700:.0f},{sub.top/12700:.0f}) {sub.width/12700:.0f}x{sub.height/12700:.0f} name='{sub.name}' text='{stxt}'")
        print(f"  Shape {i}: ({x_pt:.0f},{y_pt:.0f}) {w_pt:.0f}x{h_pt:.0f} type={shape.shape_type} name='{shape.name}' text='{txt}'")

# Get the image as base64    
img_path = r"d:\Dev\AI테스트용\pptx2web\images\shape_2.png"
with open(img_path, 'rb') as f:
    img_b64 = base64.b64encode(f.read()).decode()
print(f"\n=== Image base64 ===")
print(f"data:image/png;base64,{img_b64}")

# Also check shape 8 - note bubble text
shape8 = slide.shapes[8]
print(f"\n=== Shape 8 details ===")
print(shape8._element.xml[:2000])
