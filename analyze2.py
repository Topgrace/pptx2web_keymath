from pptx import Presentation
import sys

pptx_path = r"d:\Dev\AI테스트용\pptx2web\중등 2-1 2단원_0209(최종)_익히기합본.pptx"
prs = Presentation(pptx_path)

start = int(sys.argv[1]) - 1  # 0-indexed
end = int(sys.argv[2])  # exclusive

for slide_idx in range(start, end):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    print(f"\n{'='*60}")
    print(f"SLIDE {slide_num} ({len(slide.shapes)} shapes)")
    print(f"{'='*60}")
    
    for i, shape in enumerate(slide.shapes):
        info = f"Shape {i}: type={shape.shape_type}, name='{shape.name}'"
        
        if shape.has_text_frame:
            texts = []
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    font_info = ""
                    if para.runs:
                        r = para.runs[0]
                        font_info = f" [font={r.font.name}, sz={r.font.size}]"
                        try:
                            if r.font.color and r.font.color.rgb:
                                font_info += f" color=#{r.font.color.rgb}"
                        except: pass
                    texts.append(f"'{t}'{font_info}")
            if texts:
                print(f"  {info}")
                for t in texts:
                    print(f"    {t}")
        
        if shape.has_table:
            table = shape.table
            print(f"  {info} TABLE {len(list(table.rows))}x{len(table.columns)}")
            for r_idx, row in enumerate(table.rows):
                cells = [row.cells[c].text.strip() for c in range(len(table.columns))]
                print(f"    Row{r_idx}: {cells}")
        
        if shape.shape_type == 6:  # GROUP
            print(f"  {info} GROUP({len(shape.shapes)} children)")
            for j, child in enumerate(shape.shapes):
                if child.has_text_frame:
                    for p in child.text_frame.paragraphs:
                        if p.text.strip():
                            fi = ""
                            if p.runs:
                                r = p.runs[0]
                                fi = f" [font={r.font.name}]"
                                try:
                                    if r.font.color and r.font.color.rgb:
                                        fi += f" #{r.font.color.rgb}"
                                except: pass
                            print(f"    child{j}: '{p.text.strip()}'{fi}")
        
        if hasattr(shape, 'image'):
            try:
                img = shape.image
                print(f"  {info} IMAGE({img.content_type}, {len(img.blob)}B)")
            except: pass
