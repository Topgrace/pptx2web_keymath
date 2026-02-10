from pptx import Presentation
from pptx.oxml.ns import qn
import base64, os

pptx_path = r"d:\Dev\AI테스트용\pptx2web\중등 2-1 2단원_0209(최종)_익히기합본.pptx"
prs = Presentation(pptx_path)
slide = prs.slides[1]

# Check slide background from slide layout and master
print("=== Slide BG ===")
bg = slide.background._element
print(bg.xml[:1500] if bg is not None else "None")

# Check slide layout background
layout = slide.slide_layout
print("\n=== Layout BG ===")
lbg = layout.background._element
print(lbg.xml[:1500] if lbg is not None else "None")

# Check master background  
master = layout.slide_master
print("\n=== Master BG ===")
mbg = master.background._element
print(mbg.xml[:1500] if mbg is not None else "None")

# Get all font colors more reliably
print("\n=== All text colors ===")
for i, shape in enumerate(slide.shapes):
    if hasattr(shape, 'text_frame') and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rpr = run._r.find(qn('a:rPr'))
                color_str = ""
                if rpr is not None:
                    solid = rpr.find(qn('a:solidFill'))
                    if solid is not None:
                        srgb = solid.find(qn('a:srgbClr'))
                        if srgb is not None:
                            color_str = f"#{srgb.get('val')}"
                        scheme = solid.find(qn('a:schemeClr'))
                        if scheme is not None:
                            color_str = f"scheme:{scheme.get('val')}"
                print(f"  Shape {i}: '{run.text[:30]}' -> font={run.font.name}, size={run.font.size}, color={color_str}")

# Get shape 20 (group with math 'a') - check the fallback textbox
shape20 = slide.shapes[20]
print("\n=== Shape 20 full group XML ===")
xml = shape20._element.xml
print(xml)

# Also convert the image to base64 for embedding
img_path = r"d:\Dev\AI테스트용\pptx2web\images\shape_2.png"
with open(img_path, 'rb') as f:
    img_data = base64.b64encode(f.read()).decode()
print(f"\n=== Image base64 length: {len(img_data)} ===")
print(f"data:image/png;base64,{img_data[:100]}...")

# Check bottom formula section - the big formula aⁿ = a × a × ... × a (n번)
# Shape 5: big "2" at (186, 186)
# Shape 13: superscript "3" at (248, 174)
# These form "2³" diagram

# The bottom formula section:
# Shape arrangement for a^n = a × a × ... × a :
# Let me find all shapes in the bottom area (y > 350pt)
print("\n=== Bottom section shapes (y>350pt) ===")
for i, shape in enumerate(slide.shapes):
    y_pt = shape.top / 12700
    if y_pt > 350:
        txt = ""
        if hasattr(shape, 'text') and shape.text:
            txt = shape.text
        print(f"  Shape {i}: ({shape.left/12700:.0f},{y_pt:.0f}) {shape.width/12700:.0f}x{shape.height/12700:.0f} name='{shape.name}' text='{txt}'")
