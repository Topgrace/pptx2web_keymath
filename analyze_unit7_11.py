"""
Analyze PPTX slides for units 7-11 (II. 식의 계산)
Extracts text, fonts, colors, images, and tables from slides 27 onwards.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json, os, sys

pptx_path = r"d:\Dev\AI테스트용\pptx2web\중등 2-1 2단원_0209(최종)_익히기합본.pptx"
prs = Presentation(pptx_path)

total_slides = len(prs.slides)
print(f"Total slides in PPTX: {total_slides}")
print(f"Analyzing slides 27 to {total_slides} (1-indexed)")
print("=" * 100)

# Analyze starting from slide 27 (0-indexed: 26) to end
start_idx = 26  # 0-indexed for slide 27
end_idx = total_slides

for slide_idx in range(start_idx, end_idx):
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    print(f"\n{'='*100}")
    print(f"SLIDE {slide_num}")
    print(f"{'='*100}")

    for i, shape in enumerate(slide.shapes):
        print(f"\n--- Shape {i}: type={shape.shape_type}, name='{shape.name}' ---")
        print(f"    Position: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")

        if shape.has_text_frame:
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text
                if text.strip():
                    # Get full run-level details
                    run_details = []
                    for r in para.runs:
                        rinfo = {
                            'text': r.text,
                            'font': r.font.name,
                            'size': str(r.font.size) if r.font.size else None,
                            'bold': r.font.bold,
                            'italic': r.font.italic,
                        }
                        try:
                            if r.font.color and r.font.color.rgb:
                                rinfo['color'] = f"#{r.font.color.rgb}"
                        except:
                            pass
                        run_details.append(rinfo)

                    print(f"    P{p_idx}: '{text}'")
                    for rd in run_details:
                        print(f"        Run: '{rd['text']}' | font={rd['font']}, size={rd['size']}, bold={rd['bold']}, color={rd.get('color','?')}")

        if shape.has_table:
            table = shape.table
            print(f"    TABLE: {len(list(table.rows))} rows x {len(table.columns)} cols")
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    text = cell.text.strip()
                    if text:
                        # Get cell font colors
                        cell_colors = []
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                try:
                                    if r.font.color and r.font.color.rgb:
                                        cell_colors.append(f"#{r.font.color.rgb}")
                                except:
                                    pass
                        color_str = f" colors={cell_colors}" if cell_colors else ""
                        print(f"    Cell[{r_idx},{c_idx}]: '{text}'{color_str}")

        # Check for pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                print(f"    IMAGE: {img.content_type}, size={len(img.blob)} bytes")
            except:
                pass

        # Check if group shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"    GROUP with {len(shape.shapes)} sub-shapes")
            for j, child in enumerate(shape.shapes):
                print(f"      Child {j}: type={child.shape_type}, name='{child.name}'")
                if child.has_text_frame:
                    for p in child.text_frame.paragraphs:
                        if p.text.strip():
                            run_details = []
                            for r in p.runs:
                                rinfo = {'text': r.text, 'bold': r.font.bold}
                                try:
                                    if r.font.color and r.font.color.rgb:
                                        rinfo['color'] = f"#{r.font.color.rgb}"
                                except:
                                    pass
                                run_details.append(rinfo)
                            print(f"        '{p.text}' | runs={run_details}")
                if hasattr(child, 'image'):
                    try:
                        img = child.image
                        print(f"        IMAGE: {img.content_type}")
                    except:
                        pass

print(f"\n\n{'='*100}")
print(f"ANALYSIS COMPLETE. Processed slides {start_idx+1} to {end_idx}")
print(f"{'='*100}")
