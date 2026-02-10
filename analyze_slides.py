from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import json, os

pptx_path = r"d:\Dev\AI테스트용\pptx2web\중등 2-1 2단원_0209(최종)_익히기합본.pptx"
prs = Presentation(pptx_path)

for slide_idx in range(2, 14):  # slides 3~14 (0-indexed: 2~13)
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    print(f"\n{'='*80}")
    print(f"SLIDE {slide_num}")
    print(f"{'='*80}")
    
    for i, shape in enumerate(slide.shapes):
        print(f"\n--- Shape {i}: {shape.shape_type}, name='{shape.name}' ---")
        print(f"    Position: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}")
        
        if shape.has_text_frame:
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text
                if text.strip():
                    # Get font info from first run
                    font_info = ""
                    if para.runs:
                        r = para.runs[0]
                        font_info = f"font={r.font.name}, size={r.font.size}, bold={r.font.bold}"
                        try:
                            if r.font.color and r.font.color.rgb:
                                font_info += f", color=#{r.font.color.rgb}"
                        except: pass
                    print(f"    P{p_idx}: '{text}' [{font_info}]")
        
        if shape.has_table:
            table = shape.table
            print(f"    TABLE: {table.rows.__len__()} rows x {len(table.columns)} cols")
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    text = cell.text.strip()
                    if text:
                        print(f"    Cell[{r_idx},{c_idx}]: '{text}'")
        
        if hasattr(shape, 'image'):
            try:
                img = shape.image
                print(f"    IMAGE: {img.content_type}, size={len(img.blob)} bytes")
            except:
                pass
        
        # Check if group shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            print(f"    GROUP with {len(shape.shapes)} shapes")
            for j, child in enumerate(shape.shapes):
                print(f"      Child {j}: type={child.shape_type}, name='{child.name}'")
                if child.has_text_frame:
                    for p in child.text_frame.paragraphs:
                        if p.text.strip():
                            font_info = ""
                            if p.runs:
                                r = p.runs[0]
                                font_info = f"font={r.font.name}, size={r.font.size}, bold={r.font.bold}"
                                try:
                                    if r.font.color and r.font.color.rgb:
                                        font_info += f", color=#{r.font.color.rgb}"
                                except: pass
                            print(f"        '{p.text}' [{font_info}]")
