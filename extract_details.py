from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
import os, base64

pptx_path = r"d:\Dev\AI테스트용\pptx2web\중등 2-1 2단원_0209(최종)_익히기합본.pptx"
out_dir = r"d:\Dev\AI테스트용\pptx2web\images"
os.makedirs(out_dir, exist_ok=True)

prs = Presentation(pptx_path)
slide = prs.slides[1]

# Extract all images
for i, shape in enumerate(slide.shapes):
    if shape.shape_type == 13:  # PICTURE
        img = shape.image
        ext = img.content_type.split('/')[-1]
        if ext == 'jpeg': ext = 'jpg'
        fname = f"shape_{i}.{ext}"
        with open(os.path.join(out_dir, fname), 'wb') as f:
            f.write(img.blob)
        print(f"Saved {fname} ({len(img.blob)} bytes)")

# Extract slide background
slide_xml = slide._element.xml
print("\n=== Slide background XML snippet ===")
bg_el = slide._element.find(qn('p:bg'))
if bg_el is not None:
    print(bg_el.xml[:2000])

# Check the rounded rectangle (shape 4) for theme color
shape4 = slide.shapes[4]
print("\n=== Shape 4 (rounded rect) fill XML ===")
sp_el = shape4._element
fill_el = sp_el.find('.//' + qn('a:solidFill'))
if fill_el is not None:
    print(fill_el.xml)

# Check shape 0 line details
shape0 = slide.shapes[0]
print("\n=== Shape 0 line XML ===")
ln_el = shape0._element.find('.//' + qn('a:ln'))
if ln_el is not None:
    print(ln_el.xml[:500])

# Large "2" and "3" - the exponent diagram
# Let me check the group shapes more carefully
for i, shape in enumerate(slide.shapes):
    if shape.shape_type == 6:  # GROUP
        print(f"\n=== Group shape {i}: {shape.name} ===")
        print(shape._element.xml[:3000])

# Check freeform shape 26
print("\n=== Shape 26 (freeform) XML ===")
print(slide.shapes[26]._element.xml[:2000])

# Check the bracket in group 24
print("\n=== Shape 24 group XML ===")
print(slide.shapes[24]._element.xml[:3000])

# Also check what kind of line shape 25 (vertical arrow) is
print("\n=== Shape 25 (vertical line) ===")
print(slide.shapes[25]._element.xml[:1000])

# Extract larger number shapes content  
# Shape 5: big "2", Shape 13: "3" - they form "2^3"
# Let's also extract text boxes around formulas at the bottom

# Look at all text related to the formula section
print("\n=== Formula section shapes ===")
for i in [5, 12, 13, 18, 19, 20, 21, 22, 23, 24]:
    s = slide.shapes[i]
    print(f"Shape {i} ({s.name}): pos=({s.left/12700:.0f},{s.top/12700:.0f}) size=({s.width/12700:.0f},{s.height/12700:.0f})")
    if hasattr(s, 'text') and s.text:
        print(f"  text='{s.text}'")
